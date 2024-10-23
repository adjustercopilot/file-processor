from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from PyPDF2 import PdfReader
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
import google.generativeai as genai
from docx import Document
from pptx import Presentation
import pandas as pd
import os
import json
from dotenv import load_dotenv
from typing import List


# Load environment variables
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=api_key)

if not api_key:
    raise HTTPException(status_code=500, detail="API key not found. Please check your environment variables.")

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class ClaimFileProcessor:
    def __init__(self):
        """Initializes the generative model with Google Generative AI."""
        self.model = genai.GenerativeModel(model_name='gemini-1.5-pro')

    def get_text_from_document(self, uploaded_files: UploadFile):
        """Extract text content from various file types."""
        text = ""
        for file in uploaded_files:
            if file.filename.lower().endswith((".pdf", ".docx", ".txt", ".pptx", ".csv", ".json")):
                try:
                    if file.filename.lower().endswith(".pdf"):
                        pdf_reader = PdfReader(file.file)
                        for page in pdf_reader.pages:
                            text += page.extract_text() or ""
                    elif file.filename.lower().endswith(".docx"):
                        doc_content = Document(file.file)
                        text += "\n".join([paragraph.text for paragraph in doc_content.paragraphs])
                    elif file.filename.lower().endswith(".txt"):
                        text += file.file.read().decode('utf-8', errors='ignore')
                    elif file.filename.lower().endswith(".pptx"):
                        presentation = Presentation(file.file)
                        for slide in presentation.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + "\n"
                    elif file.filename.lower().endswith('.csv'):
                        df = pd.read_csv(file.file)
                        text += df.to_csv(index=False)
                    elif file.filename.lower().endswith('.json'):
                        json_data = json.load(file.file)
                        text  += json.dumps(json_data, indent=4) 
                except Exception as e:
                    raise HTTPException(status_code=400, detail=f"Error processing {file.filename}: {e}")
        return text
    
    def get_text_chunks(self, text):
        """Splits the extracted text into manageable chunks for processing."""
        if not text:
            raise HTTPException(status_code=400, detail="No text provided.")
        try:
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
            chunks = text_splitter.split_text(text)
            return chunks
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error splitting text: {e}")

    def extract_claim_information(self, chunks):
        """Extract specific insurance policy information from the uploaded document content."""
        prompt = """
        Extract all relevant information from the uploaded property damage document and return it in the following structured JSON format. 

        Follow these rules:
        1. If any field is missing, leave it empty.
        2. Capture all nested data accurately.
        3. Ensure numerical fields are correctly parsed (e.g., year, phone number, policy limits).
        4. For 'categories', 'subcategories', and 'questions maintain the structure as shown.

        Use the following JSON template:

        {
            "inspections": [
                {
                    "id": 1,
                    "name": "",
                    "claim_information": {
                        "claim_number": "",
                        "claim_type_id": "",
                        "created_by_id": "",
                        "customer_id": "",
                        "id": "",
                        "line_of_business_id": "",
                        "status": "",
                        "synched": 0
                    },
                    "property_information": {
                        "claim_id": "",
                        "construction_type": "",
                        "created_at": "",
                        "created_by_id": "",
                        "id": "",
                        "location_id": "",
                        "number_of_stories": 0,
                        "owner_name": "",
                        "property_type": "",
                        "roof_type": "",
                        "synched": 0,
                        "year_of_built": 0
                    },
                    "customer_information": {
                        "created_at": "",
                        "created_by_id": "",
                        "email": "",
                        "first_name": "",
                        "id": "",
                        "last_name": "",
                        "location_address": "",
                        "location_city": "",
                        "location_id": "",
                        "location_pincode": "",
                        "location_state": "",
                        "phone_number": "",
                        "updated_at": ""
                    },
                    "insurance_policy_information": {
                        "claim_id": "",
                        "coverage_type": "",
                        "created_by_id": "",
                        "deductible": "",
                        "id": "",
                        "insurance_carrier": "",
                        "policy_holder_name": "",
                        "policy_limits": "",
                        "policy_number": "",
                        "synched": 0
                    },
                    "damage_information": {
                        "cause_of_loss": "",
                        "claim_id": "",
                        "created_by_id": "",
                        "date_of_loss": "",
                        "description": "",
                        "id": "",
                        "synched": 0
                    },
                    "claim_type": {
                        "id": 1,
                        "name": ""
                    },
                    "categories": [
                        {
                        "id": 1
                        "name": "",
                        "is_interior": false,
                        "priority": 0,
                        "subcategories": [
                            {
                            "id": 1,
                            "title": "",
                            "description": "",
                            "helptext": "",
                            "priority": 1,
                            "questions": [
                                {
                                "id": 1,
                                "title": "",
                                "description": "",
                                "helptext": "",
                                "priority": 1,
                                "required": true
                                "answer_type": "",
                                "photos": true,
                                "videos": true,
                                "photos_360": true,
                                "photos_response_collection": [
                                    {}
                                ],
                                "docs": false,
                                "video_response_collection": [],
                                "360_photo_response_collection": [],
                                "notes": true,
                                "applicable": true,
                                "is_additional_questions": false
                                "response": ""
                                }
                            ]
                            }
                        ]
                        }
                    ]
                    }
                }
            ]
        }
        You do not need to adhere to any predefined format. Simply use the document structure and content to organize the output appropriately.
        Please extract relevant data from the document into this structure.
        """
        extracted_info = {}
        try:
            combined_input = f"{prompt}\n\n" + '\n'.join(chunks)
            input_data = {
                "parts": [
                    {"text": combined_input}
                ]
            }
            response = self.model.generate_content(input_data) 

            if response.candidates:
                generated_content = response.candidates[0].content.parts[0].text.strip() 
                generated_content = generated_content.replace("```json", "").replace("```", "").strip() 
                try:
                    extracted_info = json.loads(generated_content)   
                except json.JSONDecodeError as e:
                    extracted_info = {"error": "Invalid JSON response", "content": generated_content}
            else:
                extracted_info = {"error": "No candidates found in response."}

        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error during extraction: {e}")

        return extracted_info

claim_processor = ClaimFileProcessor()

@app.post("/claims/process")
async def process_claim_files(uploaded_files: List[UploadFile] = File(...)):
    """Endpoint to process uploaded claim files and extract structured policy data."""
    
    raw_text = claim_processor.get_text_from_document(uploaded_files)

    text_chunks = claim_processor.get_text_chunks(raw_text)
    
    claim_data = claim_processor.extract_claim_information(text_chunks)

    return JSONResponse(content=claim_data)
